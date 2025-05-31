import sys
print(sys.executable)

import os

import time
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv  # Add this import

# Load environment variables from .env file
load_dotenv()

# Get API key from environment variable
client = OpenAI(api_key=os.getenv("OPEN_AI"))

def extract_slide_texts(prs):
    texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
        texts.append(text.strip())
    return texts

def batch_score_slides(slide_texts, batch_size=5):
    indexed_scores = []
    for i in range(0, len(slide_texts), batch_size):
        print("Processing batch:", i)
        batch = slide_texts[i:i+batch_size]
        prompt = "Here are some presentation slides:\n\n"
        for j, text in enumerate(batch):
            prompt += f"Slide {j+1}:\n{text}\n\n"
        prompt += (
            "Please rate each slide on a scale of 1 to 10 based on importance. "
            "Only return in this format:\nSlide 1: 8\nSlide 2: 5\n..."
        )

        try:
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2,
            )
            content = response.choices[0].message.content
            for line in content.splitlines():
                if "Slide" in line and ":" in line:
                    try:
                        slide_num = int(line.split(":")[0].strip().split(" ")[-1]) - 1
                        score = int(line.split(":")[1].strip())
                        real_index = i + slide_num
                        indexed_scores.append((real_index, score))
                    except:
                        continue
        except Exception as e:
            print(f"❌ Error on batch {i}-{i+batch_size}: {e}")
            continue
        time.sleep(1.2)
    return indexed_scores

def build_trimmed_pptx(input_path, output_path, keep_indices):
    original = Presentation(input_path)
    trimmed = Presentation()
    blank_layout = trimmed.slide_layouts[6]

    for i in sorted(keep_indices):
        source_slide = original.slides[i]
        new_slide = trimmed.slides.add_slide(blank_layout)
        for shape in source_slide.shapes:
            try:
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
            except:
                continue
    trimmed.save(output_path)

def main():
    input_file = "orignal.pptx"
    output_file = "trimmed_output_15percent.pptx"

    if not os.path.exists(input_file):
        print(f"❌ File not found: {input_file}")
        return

    print("📥 Reading PPTX...")
    prs = Presentation(input_file)
    texts = extract_slide_texts(prs)

    print("🤖 Scoring slides using GPT-3.5-turbo...")
    scores = batch_score_slides(texts, batch_size=5)

    print("📊 Selecting top 15%...")
    scores.sort(key=lambda x: x[1], reverse=True)
    top_n = int(len(scores) * 0.15)
    keep_indices = [idx for idx, _ in scores[:top_n]]

    print(f"✂️ Keeping {top_n} out of {len(scores)} slides...")
    build_trimmed_pptx(input_file, output_file, keep_indices)

    print(f"✅ Done! Trimmed PPTX saved as: {output_file}")

if __name__ == "__main__":
    main()
