from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Load an existing presentation or create a new one
prs = Presentation("trimmedTemplate.pptx")  # or Presentation() for a new one

# Select the slide (or create one)
# slide = prs.slides[0]  # use an existing slide
slide = prs.slides.add_slide(prs.slide_layouts[0])  # blank slide

# Content to add
content = "This is a long or short content that will be dynamically added."

# Define position and dimensions (you can customize these)
left = Inches(1)
top = Inches(2)
width = Inches(6)
height = Inches(2)  # Can vary if needed

# Add a textbox
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Add text
p = text_frame.add_paragraph()
p.text = content

# Optional: Style the text
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)

# Save the modified presentation
prs.save("updated_presentation.pptx")
