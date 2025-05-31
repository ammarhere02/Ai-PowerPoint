import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from openai import OpenAI
from PIL import Image, ImageDraw, ImageFont
import requests
from io import BytesIO
import json
import random
from dotenv import load_dotenv
import shutil
from PIL import Image, ImageDraw, ImageFilter
# Load environment variables
load_dotenv()

class PowerPointProcessor:
    def __init__(self, openai_api_key=None, unsplash_api_key=None, template_path=None):
        """Initialize the processor with optional API keys and template path"""
        self.openai_client = None
        self.unsplash_api_key = unsplash_api_key or os.getenv("UNSPLASH_API_KEY")
        self.template_path = template_path or "trimmedTemplate.pptx"
        
        # Use environment variable if no key provided
        openai_api_key = openai_api_key or os.getenv("OPEN_AI")
        if openai_api_key:
            self.openai_client = OpenAI(api_key=openai_api_key)
        
    def extract_text_from_pptx(self, file_path):
        """Extract text content from PowerPoint slides"""
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            slides_content.append({
                'slide_number': i + 1,
                'content': slide_text
            })
        
        return slides_content
    
    def _extract_slide_texts(self, prs):
        """Extract text content from PowerPoint slides with better handling of long text"""
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Append text with a newline separator
                    slide_text += shape.text.strip() + "\n\n"
            texts.append({
                'slide_number': i + 1,
                'content': slide_text.strip()
            })
        return texts
    
    def structure_content(self, slides_content):
        """Convert slides to properly structured format with headings, tables, and bullet points"""
        structured_slides = []
        
        for slide in slides_content:
            slide_num = slide['slide_number']
            content = slide['content']
            
            structured_content = self._analyze_and_structure_text(content)
            
            structured_slides.append({
                'slide_number': slide_num,
                'structured_content': structured_content,
                'needs_image': self._determine_image_need(structured_content)
            })
        
        return structured_slides
    
    def _analyze_and_structure_text(self, content):
        """Analyze text content and structure it with appropriate formatting"""
        if not content:
            return {'type': 'empty', 'title': 'Empty Slide', 'data': []}
        
        combined_text = '\n'.join(content)
        
        # Detect if content looks like tabular data
        if self._is_tabular_data(combined_text):
            return {
                'type': 'table',
                'title': self._extract_title(content[0] if content else 'Data Table'),
                'data': self._parse_table_data(combined_text)
            }
        
        # Detect if content has bullet points or lists
        elif self._has_bullet_points(combined_text):
            return {
                'type': 'bullet_list',
                'title': self._extract_title(content[0] if content else 'Key Points'),
                'points': self._extract_bullet_points(combined_text)
            }
        
        # Default to structured text with heading
        else:
            return {
                'type': 'structured_text',
                'title': self._extract_title(content[0] if content else 'Content'),
                'content': self._structure_paragraphs(combined_text)
            }
    
    def _is_tabular_data(self, text):
        """Check if text contains tabular data patterns"""
        patterns = [
            r'\t.*\t',  # Tab-separated
            r'\|.*\|',  # Pipe-separated
            r':\s*\d+',  # Colon-number pairs
            r'\d+\.\d+%',  # Percentages
        ]
        return any(re.search(pattern, text) for pattern in patterns)
    
    def _has_bullet_points(self, text):
        """Check if text has bullet point patterns"""
        patterns = [
            r'^[\•\-\*]\s',  # Bullet symbols
            r'^\d+\.\s',     # Numbered lists
            r'^[a-zA-Z]\.\s' # Letter lists
        ]
        lines = text.split('\n')
        if not lines:
            return False
        bullet_lines = sum(1 for line in lines if any(re.match(pattern, line.strip()) for pattern in patterns))
        return bullet_lines > len(lines) * 0.3  # 30% of lines are bullet points
    
    def _extract_title(self, first_line):
        """Extract title from first line"""
        if not first_line:
            return "Slide Content"
        if len(first_line) < 80 and not first_line.endswith('.'):
            return first_line
        return "Key Points"
    
    def _parse_table_data(self, text):
        """Parse text into table format"""
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        table_data = []
        
        for line in lines[1:]:  # Skip title line
            # Try different separators
            if '\t' in line:
                row = [cell.strip() for cell in line.split('\t')]
            elif '|' in line:
                row = [cell.strip() for cell in line.split('|')]
            else:
                # Try to split on multiple spaces or colons
                row = re.split(r'\s{2,}|:', line)
                row = [cell.strip() for cell in row if cell.strip()]
            
            if len(row) > 1:
                table_data.append(row)
        
        return table_data
    
    def _extract_bullet_points(self, text):
        """Extract bullet points from text"""
        lines = text.split('\n')
        points = []
        
        for line in lines:
            line = line.strip()
            if re.match(r'^[\•\-\*]\s', line):
                points.append(re.sub(r'^[\•\-\*]\s', '', line))
            elif re.match(r'^\d+\.\s', line):
                points.append(re.sub(r'^\d+\.\s', '', line))
            elif re.match(r'^[a-zA-Z]\.\s', line):
                points.append(re.sub(r'^[a-zA-Z]\.\s', '', line))
            elif line and not any(re.match(r'^[\•\-\*]\s|^\d+\.\s|^[a-zA-Z]\.\s', l.strip()) for l in lines[:max(0, lines.index(line))]):
                # This might be a title or header - skip
                continue
            elif line:
                points.append(line)
        
        return points if points else ["Content point"]
    
    def _structure_paragraphs(self, text):
        """Structure text into paragraphs with better handling of long text"""
        # First try to split by double newlines (paragraph breaks)
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        # If no paragraph breaks found, try single newlines
        if not paragraphs:
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        # If still no breaks, handle as a single paragraph
        if not paragraphs:
            paragraphs = [text]
        
        # Ensure we don't lose any content by truncating
        return paragraphs
    
    def _determine_image_need(self, structured_content):
        """Determine if a slide would benefit from an image"""
        image_keywords = [
            'process', 'workflow', 'diagram', 'chart', 'graph', 'visual',
            'example', 'comparison', 'analysis', 'data', 'statistics',
            'architecture', 'design', 'model', 'framework', 'structure',
            'timeline', 'roadmap', 'overview', 'summary', 'business',
            'strategy', 'growth', 'performance', 'results', 'metrics'
        ]
        
        content_text = str(structured_content).lower()
        keyword_matches = sum(1 for keyword in image_keywords if keyword in content_text)
        
        # Also consider content type
        if structured_content['type'] in ['table', 'bullet_list']:
            return True
        
        return keyword_matches >= 2
    
    def generate_images_for_slides(self, structured_slides):
        """Generate images for slides that need them"""
        for slide in structured_slides:
            if slide['needs_image']:
                image_prompt = self._create_image_prompt(slide['structured_content'])
                image_path = self._generate_image(image_prompt, slide['slide_number'])
                slide['image_path'] = image_path
            else:
                slide['image_path'] = None
        
        return structured_slides
    
    def _create_image_prompt(self, content):
        """Create a prompt for image generation based on content"""
        content_type = content['type']
        title = content.get('title', 'Slide Content')
        
        if content_type == 'table':
            return "business data analytics dashboard charts graphs professional"
        elif content_type == 'bullet_list':
            return "business presentation infographic professional corporate strategy"
        else:
            return "business professional office meeting presentation corporate"
    
    def _generate_image(self, prompt, slide_number):
        """Generate image using available APIs or create a placeholder"""
        # Create pics directory if it doesn't exist
        pics_dir = "pics"
        if not os.path.exists(pics_dir):
            os.makedirs(pics_dir)
        
        try:
            # Try Unsplash first
            if self.unsplash_api_key:
                return self._get_unsplash_image(prompt, slide_number, pics_dir)
            
            # Try OpenAI DALL-E (new API)
            elif self.openai_client:
                return self._generate_openai_image(prompt, slide_number, pics_dir)
            
            else:
                # Create placeholder image
                return self._create_placeholder_image(prompt, slide_number, pics_dir)
    
        except Exception as e:
            print(f"Error generating image for slide {slide_number}: {e}")
            return self._create_placeholder_image(prompt, slide_number, pics_dir)

    def _get_unsplash_image(self, prompt, slide_number, pics_dir):
        """Get image from Unsplash API"""
        try:
            url = "https://api.unsplash.com/search/photos"
            headers = {"Authorization": f"Client-ID {self.unsplash_api_key}"}
            params = {
                "query": prompt,
                "per_page": 10,
                "orientation": "landscape"
            }
            
            response = requests.get(url, headers=headers, params=params)
            response.raise_for_status()
            
            data = response.json()
            
            if data['results']:
                # Get a random image from results
                image_data = random.choice(data['results'])
                image_url = image_data['urls']['regular']
                
                # Download and save image
                img_response = requests.get(image_url)
                img_response.raise_for_status()
                
                image = Image.open(BytesIO(img_response.content))
                
                # Resize image for presentation
                image = image.resize((800, 600), Image.Resampling.LANCZOS)
                
                image_path = os.path.join(pics_dir, f"unsplash_image_slide_{slide_number}.jpg")
                image.save(image_path, "JPEG", quality=85)
                return image_path
            
            else:
                return self._create_placeholder_image(prompt, slide_number, pics_dir)
                
        except Exception as e:
            print(f"Unsplash API error: {e}")
            return self._create_placeholder_image(prompt, slide_number, pics_dir)

    def _generate_openai_image(self, prompt, slide_number, pics_dir):
        """Generate image using OpenAI DALL-E (new API)"""
        try:
            full_prompt = f"Create a professional business presentation image for: {prompt}. Modern, clean design suitable for corporate presentation. High quality, professional photography style."
            
            response = self.openai_client.images.generate(
                model="dall-e-3",
                prompt=full_prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            
            image_url = response.data[0].url
            
            # Download and save image
            img_response = requests.get(image_url)
            img_response.raise_for_status()
            
            image = Image.open(BytesIO(img_response.content))
            
            # Resize for presentation
            image = image.resize((800, 600), Image.Resampling.LANCZOS)
            
            image_path = os.path.join(pics_dir, f"openai_image_slide_{slide_number}.jpg")
            image.save(image_path, "JPEG", quality=85)
            return image_path
            
        except Exception as e:
            print(f"OpenAI API error: {e}")
            return self._create_placeholder_image(prompt, slide_number, pics_dir)

    def _create_placeholder_image(self, prompt, slide_number, pics_dir):
        """Create a placeholder image when APIs are not available"""
        img = Image.new('RGB', (800, 600), color='#f8f9fa')
        draw = ImageDraw.Draw(img)
        
        # Draw a simple border
        draw.rectangle([10, 10, 790, 590], outline='#dee2e6', width=2)
        
        try:
            font_large = ImageFont.truetype("arial.ttf", 24)
            font_small = ImageFont.truetype("arial.ttf", 16)
        except:
            font_large = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        # Add title
        title_text = f"Slide {slide_number} Image"
        
        # Calculate text position for centering
        bbox = draw.textbbox((0, 0), title_text, font=font_large)
        text_width = bbox[2] - bbox[0]
        x = (800 - text_width) // 2
        
        draw.text((x, 250), title_text, fill='#495057', font=font_large)
        
        # Add description
        desc_text = f"Image placeholder for: {prompt[:50]}..."
        bbox = draw.textbbox((0, 0), desc_text, font=font_small)
        text_width = bbox[2] - bbox[0]
        x = (800 - text_width) // 2
        
        draw.text((x, 300), desc_text, fill='#6c757d', font=font_small)
        
        # Add simple graphic element
        draw.ellipse([350, 350, 450, 450], fill='#007bff', outline='#0056b3', width=2)
        
        image_path = os.path.join(pics_dir, f"placeholder_image_slide_{slide_number}.png")
        img.save(image_path)
        return image_path
    
    def create_presentation_from_template(self, structured_slides, output_path="enhanced_presentation.pptx"):
        """Create presentation using existing template"""
        try:
            # Load the existing template
            if os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                print(f"✅ Loaded template: {self.template_path}")
            else:
                print(f"⚠️ Template not found: {self.template_path}, creating new presentation")
                prs = Presentation()
        except Exception as e:
            print(f"⚠️ Error loading template: {e}, creating new presentation")
            prs = Presentation()
        
        # Clear existing slides except master/layout slides
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            if i < len(prs.slides):
                try:
                    xml_slides = prs.slides._sldIdLst
                    xml_slides.remove(xml_slides[i])
                except:
                    pass
        
        # Get available layouts
        available_layouts = len(prs.slide_layouts)
        print(f"📝 Available slide layouts: {available_layouts}")
        
        # Add slides using template layouts
        for slide_data in structured_slides:
            # Choose appropriate layout (try layout 1 for content, fallback to 0)
            try:
                if available_layouts > 1:
                    layout = prs.slide_layouts[1]  # Usually content layout
                else:
                    layout = prs.slide_layouts[0]  # Title or blank layout
            except:
                layout = prs.slide_layouts[0]  # Fallback to first available
            
            slide = prs.slides.add_slide(layout)
            self._populate_template_slide(slide, slide_data)
        
        # Save the presentation
        prs.save(output_path)
        return output_path
    
    def _populate_template_slide(self, slide, slide_data):
        """Populate a template slide with content"""
        content = slide_data['structured_content']
        image_path = slide_data.get('image_path')
        
        # Try to use existing placeholders first
        title_filled = False
        content_filled = False
        
        # Fill title placeholder if available
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    # This is a placeholder
                    if hasattr(shape, 'text_frame') and not title_filled:
                        # Title placeholder
                        if shape.text_frame is not None:
                            shape.text_frame.text = content['title']
                            # Style the title
                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.font.size = Pt(24)
                                paragraph.font.bold = True
                                paragraph.alignment = PP_ALIGN.CENTER  # Center align title
                            title_filled = True
                    elif hasattr(shape, 'text_frame') and not content_filled and title_filled:
                        # Content placeholder
                        self._fill_content_placeholder(shape, content)
                        content_filled = True
        
        # If no placeholders found or not filled, add textboxes manually
        if not title_filled:
            self._add_title_textbox(slide, content['title'])
        
        if not content_filled:
            self._add_content_textbox(slide, content, image_path)
        
        # Add image if available and there's space
        if image_path and os.path.exists(image_path):
            self._add_image_to_slide(slide, image_path)
    
    def _fill_content_placeholder(self, shape, content):
        """Fill content placeholder based on content type with centered alignment"""
        text_frame = shape.text_frame
        text_frame.clear()  # Clear existing content
        
        # Set text frame properties for better centering
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
        
        if content['type'] == 'bullet_list':
            points = content.get('points', [])
            for i, point in enumerate(points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Process formatting in the text
                formatted_parts = self._process_text_formatting(point)
                if len(formatted_parts) == 1 and not formatted_parts[0][1]:
                    # No special formatting
                    p.text = formatted_parts[0][0]
                else:
                    # Apply formatting
                    p.text = ""
                    for text_part, is_bold in formatted_parts:
                        run = p.add_run()
                        run.text = text_part
                        if is_bold:
                            run.font.bold = True
                
                p.level = 0  # First level bullet
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.CENTER  # Center align each bullet point
                p.space_after = Pt(12)  # Add space between points
        
        elif content['type'] == 'structured_text':
            paragraphs = content.get('content', [])
            for i, paragraph in enumerate(paragraphs):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = str(paragraph)
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.CENTER  # Center align paragraphs
                p.space_after = Pt(12)
        
        else:  # table or other content
            # For tables, convert to bullet points or simple text
            if content['type'] == 'table' and content.get('data'):
                for i, row in enumerate(content['data'][:5]):  # Limit rows
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = " | ".join(str(cell) for cell in row)
                    p.font.size = Pt(14)
                    p.alignment = PP_ALIGN.CENTER  # Center align table rows
                    p.space_after = Pt(8)
            else:
                p = text_frame.paragraphs[0]
                p.text = "Content available"
                p.alignment = PP_ALIGN.CENTER
                
    def _add_title_textbox(self, slide, title):
        """Add title textbox if no title placeholder exists with center alignment"""
        left = Inches(2)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1.2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Process formatting in the title
        formatted_parts = self._process_text_formatting(title)
        if len(formatted_parts) == 1 and not formatted_parts[0][1]:
            # No special formatting
            text_frame.text = formatted_parts[0][0]
            p = text_frame.paragraphs[0]
        else:
            # Apply formatting
            p = text_frame.paragraphs[0]
            p.text = ""
            for text_part, is_bold in formatted_parts:
                run = p.add_run()
                run.text = text_part
                if is_bold:
                    run.font.bold = True
        
        # Set text frame properties
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Style the title
        p.font.size = Pt(28)
        p.font.bold = True  # The entire title is already bold by default
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER  # Center align title
    
    def _add_content_textbox(self, slide, content, image_path):
        """Add content textbox with centered alignment and appropriate layout"""
        # Adjust layout based on whether image is present
        if image_path and os.path.exists(image_path):
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(5.5)
            height = Inches(4)
        else:
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Set text frame properties for better centering
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
        
        # Populate content based on type
        if content['type'] == 'bullet_list':
            points = content.get('points', [])
            for i, point in enumerate(points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                    
                # Process formatting in the text
                formatted_parts = self._process_text_formatting(point)
                if len(formatted_parts) == 1 and not formatted_parts[0][1]:
                    # No special formatting, add bullet symbol
                    p.text = f"• {formatted_parts[0][0]}"
                else:
                    # Apply formatting with bullet
                    p.text = "• "  # Start with bullet
                    for text_part, is_bold in formatted_parts:
                        run = p.add_run()
                        run.text = text_part
                        if is_bold:
                            run.font.bold = True
                            
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER  # Center align bullet points
                p.space_after = Pt(12)
        
        elif content['type'] == 'structured_text':
            paragraphs = content.get('content', [])
            for i, paragraph in enumerate(paragraphs):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                    
                # Process formatting in the text
                formatted_parts = self._process_text_formatting(str(paragraph))
                if len(formatted_parts) == 1 and not formatted_parts[0][1]:
                    # No special formatting
                    p.text = formatted_parts[0][0]
                else:
                    # Apply formatting
                    p.text = ""
                    for text_part, is_bold in formatted_parts:
                        run = p.add_run()
                        run.text = text_part
                        if is_bold:
                            run.font.bold = True
                            
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER  # Center align paragraphs
                p.space_after = Pt(12)
        
        elif content['type'] == 'table':
            # Convert table to formatted text
            table_data = content.get('data', [])
            for i, row in enumerate(table_data[:8]):  # Limit rows
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                    
                # Join the cells but process each for formatting
                row_text = " | ".join(str(cell) for cell in row)
                formatted_parts = self._process_text_formatting(row_text)
                
                if len(formatted_parts) == 1 and not formatted_parts[0][1]:
                    # No special formatting
                    p.text = formatted_parts[0][0]
                else:
                    # Apply formatting
                    p.text = ""
                    for text_part, is_bold in formatted_parts:
                        run = p.add_run()
                        run.text = text_part
                        if is_bold:
                            run.font.bold = True
                            
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER  # Center align table rows
                p.space_after = Pt(8)

    def add_fade_to_edges(self,image_path, output_path):
        img = Image.open(image_path).convert("RGBA")

        # Create gradient mask
        mask = Image.new("L", img.size, 255)
        draw = ImageDraw.Draw(mask)

        # Feather/fade area size
        fade_width = int(min(img.size) * 0.15)

        # Apply a radial gradient for fading
        for x in range(img.size[0]):
            for y in range(img.size[1]):
                distance_to_edge = min(x, y, img.size[0]-x, img.size[1]-y)
                if distance_to_edge < fade_width:
                    fade = int(255 * (distance_to_edge / fade_width))
                    mask.putpixel((x, y), fade)

        img.putalpha(mask)
        img.save(output_path, "PNG")
        
    def _add_image_to_slide(self, slide, image_path):
        """Add image to slide with appropriate positioning"""
        
        self.add_fade_to_edges(image_path,f'{image_path}.faded')
        try:
            left = Inches(7.5)
            top = Inches(2.5)
            width = Inches(3)
            height = Inches(3)
            
            slide.shapes.add_picture(f'{image_path}.faded', left, top, width, height)
        except Exception as e:
            print(f"Error adding image to slide: {e}")
    def _process_text_formatting(self, text):
        """Process text to apply formatting patterns like **bold**"""
        # Check if the text contains bold pattern
        if "**" in text:
            # Extract text between ** markers
            parts = []
            is_bold = False
            current_part = ""
            
            i = 0
            while i < len(text):
                if i + 1 < len(text) and text[i:i+2] == "**":
                    # Found a bold marker, toggle bold state
                    parts.append((current_part, is_bold))
                    current_part = ""
                    is_bold = not is_bold
                    i += 2
                else:
                    current_part += text[i]
                    i += 1
                    
            # Add the last part
            if current_part:
                parts.append((current_part, is_bold))
                
            return parts
        else:
            # No special formatting, return as regular text
            return [(text, False)]
    def process_presentation(self, input_pptx_path, output_pptx_path="enhanced_presentation.pptx"):
        """Main method to process the entire presentation using template"""
        print("Step 1: Extracting content from PowerPoint...")
        slides_content = self.extract_text_from_pptx(input_pptx_path)
        
        print("Step 2: Structuring content...")
        structured_slides = self.structure_content(slides_content)
        
        print("Step 3: Generating images for relevant slides...")
        slides_with_images = self.generate_images_for_slides(structured_slides)
        
        print("Step 4: Creating presentation from template...")
        final_presentation_path = self.create_presentation_from_template(slides_with_images, output_pptx_path)
        
        print(f"✅ Enhanced presentation created: {final_presentation_path}")
        
        # Print summary
        image_slides = sum(1 for slide in slides_with_images if slide['needs_image'])
        print(f"📊 Processed {len(slides_content)} slides")
        print(f"🖼️  Generated images for {image_slides} slides")
        print(f"🎨 Used template: {self.template_path}")
        
        # Delete the pics folder after processing is complete
        if os.path.exists("pics"):
            try:
                shutil.rmtree("pics")
                print("🗑️  Temporary image files cleaned up")
            except Exception as e:
                print(f"⚠️ Warning: Could not remove temporary images: {e}")
    
        return final_presentation_path


# Usage example
if __name__ == "__main__":
    # Initialize processor with API keys and template path
    processor = PowerPointProcessor(template_path="trimmedTemplate.pptx")
    
    # Process presentation
    input_file = "trimmed_output_15percent.pptx"  # Your input file
    
    output_file = "enhanced_presentation.pptx"
    
    if os.path.exists(input_file):
        try:
            result = processor.process_presentation(input_file, output_file)
            print(f"\n🎉 Success! Enhanced presentation saved as: {result}")
        except Exception as e:
            print(f"❌ Error processing presentation: {e}")
            import traceback
            traceback.print_exc()
    else:
        print(f"❌ Input file '{input_file}' not found!")
        print("Creating a sample presentation for testing...")
        
        # Create a sample presentation for testing
        sample_prs = Presentation()
        slide = sample_prs.slides[0]
        slide.shapes.title.text = "Sample Business Report"
        content = slide.placeholders[1].text_frame
        content.text = "Key Performance Indicators:\n• Revenue increased by 25%\n• Customer satisfaction: 92%\n• Market share growth: 15%"
        
        sample_prs.save("sample_presentation.pptx")
        print("Sample presentation created: sample_presentation.pptx")
        
        # Process the sample
        try:
            result = processor.process_presentation("sample_presentation.pptx", "enhanced_sample.pptx")
            print(f"\n🎉 Sample processed successfully: {result}")
        except Exception as e:
            print(f"❌ Error processing sample: {e}")
            import traceback
            traceback.print_exc()